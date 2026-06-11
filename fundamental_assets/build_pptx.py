"""Render a Fundamental HTML deck to PPTX as one full-bleed image per slide.

Usage:
    python fundamental_assets/build_pptx.py module_00
    python fundamental_assets/build_pptx.py module_00/module_00_presentation.html

Each slide is captured at 1920x1080 via headless Chrome (?slide=N&noscale=1),
assembled into a 13.33"x7.5" widescreen deck, and speaker notes from the HTML's
`<script id="speaker-notes">` JSON block are attached by index.

The HTML remains the source of truth. Re-run this script after edits.
"""
from __future__ import annotations

import json
import re
import shutil
import subprocess
import sys
import tempfile
import time
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
SLIDE_W, SLIDE_H = 1920, 1080
# Chrome --headless=new on macOS under-reports the usable viewport height
# relative to --window-size (observed: a 1920x1080 window yields only 993
# CSS pixels of rendered content). We give Chrome extra vertical slack and
# crop the PNG back to the authored size after capture.
CAPTURE_H_SLACK = 200
DPR = 2

# PowerPoint canonical widescreen 16:9 in EMU (1 inch = 914400 EMU).
# Using exact EMU values instead of Inches(13.333) prevents a ~305 EMU
# rounding shortfall that makes viewers (Keynote, LibreOffice) render
# white gaps when they snap the slide to standard 16:9 on import.
SLIDE_W_EMU = 12192000  # 13.3333... in
SLIDE_H_EMU = 6858000   # 7.5 in
# Tiny overlap per side so any sub-pixel rounding in viewers crops
# imperceptibly instead of exposing paper background.
BLEED_EMU = 9144        # 0.01 in


def resolve_html(arg: str) -> Path:
    p = Path(arg).resolve()
    if p.is_dir():
        candidates = sorted(p.glob("*_presentation.html"))
        if not candidates:
            raise SystemExit(f"no *_presentation.html under {p}")
        return candidates[0]
    if p.suffix != ".html":
        raise SystemExit(f"expected an HTML file or directory, got {p}")
    return p


def extract_speaker_notes(html: Path) -> list[str]:
    text = html.read_text(encoding="utf-8")
    m = re.search(
        r'<script type="application/json" id="speaker-notes">\s*(\[.*?\])\s*</script>',
        text,
        re.DOTALL,
    )
    if not m:
        raise SystemExit(f"no speaker-notes JSON block in {html}")
    return json.loads(m.group(1))


def count_slides(html: Path) -> int:
    text = html.read_text(encoding="utf-8")
    return len(re.findall(r"<section\s[^>]*data-screen-label=", text))


def capture_slide(html: Path, index: int, out_png: Path, profile_root: Path) -> None:
    """Screenshot one slide. Chrome doesn't exit cleanly after --screenshot,
    so we poll for the output file and terminate the process once it's stable."""
    profile = profile_root / f"p{index:02d}"
    profile.mkdir(parents=True, exist_ok=True)

    url = f"file://{html}?slide={index}&noscale=1"
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--hide-scrollbars",
        "--no-first-run",
        "--no-default-browser-check",
        "--disable-background-networking",
        "--disable-default-apps",
        "--disable-extensions",
        "--disable-sync",
        f"--user-data-dir={profile}",
        f"--window-size={SLIDE_W},{SLIDE_H + CAPTURE_H_SLACK}",
        f"--force-device-scale-factor={DPR}",
        "--virtual-time-budget=5000",
        f"--screenshot={out_png}",
        url,
    ]
    if out_png.exists():
        out_png.unlink()

    proc = subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    deadline = time.time() + 30
    last_size, stable = -1, 0
    try:
        while time.time() < deadline:
            if out_png.exists():
                size = out_png.stat().st_size
                if size > 1000 and size == last_size:
                    stable += 1
                    if stable >= 2:
                        break
                else:
                    stable = 0
                last_size = size
            time.sleep(0.2)
        else:
            raise SystemExit(f"capture timed out for slide {index} of {html.name}")
    finally:
        proc.terminate()
        try:
            proc.wait(timeout=2)
        except subprocess.TimeoutExpired:
            proc.kill()

    # Crop to authored size (the bottom slack was captured as white padding
    # because Chrome's effective viewport is shorter than --window-size).
    _crop_to_authored(out_png)


def _crop_to_authored(png_path: Path) -> None:
    """Crop the captured PNG to SLIDE_W x SLIDE_H at DPR — drops the extra
    vertical slack we gave Chrome."""
    from PIL import Image

    target_w, target_h = SLIDE_W * DPR, SLIDE_H * DPR
    with Image.open(png_path) as im:
        w, h = im.size
        if w == target_w and h == target_h:
            return  # already correct size
        # Crop from top-left; the authored content starts at (0, 0) and the
        # slack is always at the bottom.
        cropped = im.crop((0, 0, target_w, target_h))
        cropped.save(png_path, "PNG")


def _strip_placeholders(prs: Presentation) -> None:
    """Remove all placeholder shapes from the slide master and every layout.

    python-pptx's default template ships with a 4:3 master whose placeholders
    (Title, Text, Date, Footer, Slide Number) sit at positions laid out for
    a 10"x7.5" canvas. When we resize the presentation to 16:9, those
    placeholders remain at their 4:3 coordinates — and PowerPoint appears
    to reserve visual space for the bottom row (Date/Footer/Slide Number
    at y=6.95"), which produces the ~10% white band at the bottom of
    full-bleed picture slides.
    Stripping them leaves a truly blank master/layout chain so the slide's
    picture renders edge-to-edge."""
    spTree_tag = qn("p:spTree")
    sp_tag = qn("p:sp")
    containers = [prs.slide_master.element] + [l.element for l in prs.slide_layouts]
    for el in containers:
        spTree = el.find(f".//{spTree_tag}")
        if spTree is None:
            continue
        for sp in list(spTree.findall(sp_tag)):
            spTree.remove(sp)


def _set_slide_background_image(slide, image_path: Path) -> None:
    """Attach the image to the slide as its <p:bg>, not as a picture shape.

    Picture shapes can be clipped by PowerPoint's slide-area layout logic
    (observed empirically on macOS PowerPoint: the bottom ~10% of a
    picture shape was dropped despite the shape being correctly sized
    and positioned). Slide backgrounds via <p:bg> are rendered first and
    always fill the slide area edge-to-edge.

    Approach: use add_picture() to get the image embedded as a related
    part with a valid rId, read that rId, delete the picture shape, and
    emit a <p:bg><p:bgPr><a:blipFill ...> referencing the same rId."""
    from lxml import etree

    # Use add_picture as a convenient way to register the image and get an rId.
    pic = slide.shapes.add_picture(str(image_path), left=Emu(0), top=Emu(0))
    pic_el = pic._element
    blip = pic_el.find(f".//{qn('a:blip')}")
    r_embed_attr = qn("r:embed")
    rel_id = blip.get(r_embed_attr)
    # Remove the picture shape — only the related image part + rId remain.
    pic_el.getparent().remove(pic_el)

    # Build <p:bg> with a blipFill that stretches the image to the full slide.
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    bg_xml = (
        f'<p:bg xmlns:p="{ns_p}" xmlns:a="{ns_a}" xmlns:r="{ns_r}">'
        f'<p:bgPr><a:blipFill dpi="0" rotWithShape="1">'
        f'<a:blip r:embed="{rel_id}"/>'
        f'<a:srcRect/><a:stretch><a:fillRect/></a:stretch>'
        f'</a:blipFill><p:effectLst/></p:bgPr></p:bg>'
    )
    bg_el = etree.fromstring(bg_xml)

    cSld = slide.element.find(qn("p:cSld"))
    spTree = cSld.find(qn("p:spTree"))
    # <p:bg> must precede <p:spTree> per schema.
    cSld.insert(list(cSld).index(spTree), bg_el)


def build_pptx(pngs: list[Path], notes: list[str], out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    # python-pptx inherits type="screen4x3" from its default template, which
    # makes some viewers (Keynote especially) render at 4:3 and crop the
    # bottom of full-bleed 16:9 pictures. Force the type to match.
    sld_sz = prs._element.find(qn("p:sldSz"))
    sld_sz.set("type", "screen16x9")
    # Remove master+layout placeholders so PowerPoint does not reserve
    # bottom chrome space for empty Date/Footer/Slide Number fields.
    _strip_placeholders(prs)
    blank = prs.slide_layouts[6]
    for png, note in zip(pngs, notes):
        s = prs.slides.add_slide(blank)
        _set_slide_background_image(s, png)
        s.notes_slide.notes_text_frame.text = note
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main(argv: list[str]) -> int:
    if len(argv) != 2:
        print(__doc__)
        return 2

    html = resolve_html(argv[1])
    notes = extract_speaker_notes(html)
    n = count_slides(html)
    if len(notes) != n:
        raise SystemExit(
            f"{html.name}: {n} slides but {len(notes)} speaker notes — they must match"
        )

    out_pptx = html.with_suffix(".pptx")
    print(f"rendering {n} slides from {html}")

    with tempfile.TemporaryDirectory(prefix="fd-capture-") as tmp:
        tmp = Path(tmp)
        pngs: list[Path] = []
        for i in range(n):
            png = tmp / f"slide_{i:02d}.png"
            print(f"  [{i + 1:2d}/{n}] {png.name}")
            capture_slide(html, i, png, tmp / "profiles")
            pngs.append(png)

        # Stash a copy for debugging
        debug = Path("/tmp/fundamental_slides")
        debug.mkdir(exist_ok=True)
        for p in pngs:
            shutil.copy(p, debug / p.name)

        build_pptx(pngs, notes, out_pptx)

    print(f"done -> {out_pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
